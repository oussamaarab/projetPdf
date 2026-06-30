#!/usr/bin/env python3
"""
pdf_to_pptx.py – Convert a directory of page images into a PPTX file.

Usage:
    python3 pdf_to_pptx.py <images_directory> <output_pptx_path>

Works on Linux (Docker / Render) and Windows without modification.
"""

import glob
import os
import platform
import re
import site
import sys


# ── Ensure user site-packages are on sys.path (Windows only) ─────────────────
# On Linux the packages are installed system-wide via pip3 (see Dockerfile),
# so no path manipulation is needed.
def _add_windows_user_site_packages() -> None:
    """Probe common Windows user site-packages locations and add them to sys.path."""
    localappdata = os.environ.get("LOCALAPPDATA", "")
    appdata      = os.environ.get("APPDATA",      "")

    candidates = []

    # Windows Store Python (3.11 / 3.12)
    store_base = os.path.join(localappdata, "Packages")
    if os.path.isdir(store_base):
        for entry in os.scandir(store_base):
            if entry.is_dir() and "PythonSoftwareFoundation" in entry.name:
                for ver in ("Python311", "Python312", "Python313"):
                    candidates.append(
                        os.path.join(
                            entry.path,
                            "LocalCache", "local-packages", ver, "site-packages",
                        )
                    )

    # Standard user installations (%APPDATA%\Python\PythonXYZ\site-packages)
    if appdata:
        for ver in ("Python311", "Python312", "Python313"):
            candidates.append(os.path.join(appdata, "Python", ver, "site-packages"))

    # Standard system installations (%LOCALAPPDATA%\Programs\Python)
    if localappdata:
        prog_python = os.path.join(localappdata, "Programs", "Python")
        if os.path.isdir(prog_python):
            for entry in os.scandir(prog_python):
                if entry.is_dir():
                    candidates.append(
                        os.path.join(entry.path, "Lib", "site-packages")
                    )

    # Standard site.getusersitepackages() (works on all platforms)
    try:
        candidates.append(site.getusersitepackages())
    except AttributeError:
        pass

    for path in candidates:
        if os.path.isdir(path) and path not in sys.path:
            sys.path.append(path)


if platform.system() == "Windows":
    _add_windows_user_site_packages()


# ── Imports that require python-pptx / Pillow ─────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
except ImportError as exc:
    print(
        f"ERROR: Required Python package not found: {exc}\n"
        "Install it with:  pip3 install python-pptx Pillow",
        file=sys.stderr,
    )
    sys.exit(1)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _natural_sort_key(s: str):
    """Sort strings with embedded numbers in human order (page-2 < page-10)."""
    return [
        int(part) if part.isdigit() else part.lower()
        for part in re.split(r"(\d+)", s)
    ]


# ── Main conversion function ──────────────────────────────────────────────────

def pdf_pages_to_pptx(images_dir: str, output_pptx: str) -> None:
    """Package all PNG/JPG page images in *images_dir* into a PPTX file."""

    # Collect images regardless of OS path separator
    images: list[str] = []
    for ext in ("*.png", "*.jpg", "*.jpeg"):
        images.extend(glob.glob(os.path.join(images_dir, ext)))

    images.sort(key=_natural_sort_key)

    if not images:
        print(f"ERROR: No page images found in '{images_dir}'.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank

    # Set slide dimensions from the first image (assumes 150 DPI render)
    DPI = 150
    with Image.open(images[0]) as first:
        img_w, img_h = first.size

    prs.slide_width  = Inches(img_w / DPI)
    prs.slide_height = Inches(img_h / DPI)

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, 0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Ensure the output directory exists
    os.makedirs(os.path.dirname(os.path.abspath(output_pptx)), exist_ok=True)

    prs.save(output_pptx)
    print(f"Presentation saved to: {output_pptx}")


# ── Entry point ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(
            "Usage: python3 pdf_to_pptx.py <images_directory> <output_pptx_path>",
            file=sys.stderr,
        )
        sys.exit(1)

    pdf_pages_to_pptx(sys.argv[1], sys.argv[2])
